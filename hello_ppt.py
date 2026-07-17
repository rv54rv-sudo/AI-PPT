from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor

# 建立簡報
prs = Presentation()

# 設定 16:9
prs.slide_width = Cm(33.867)
prs.slide_height = Cm(19.05)

slide = prs.slides.add_slide(prs.slide_layouts[6])

# ===== 標題 =====
textbox = slide.shapes.add_textbox(
    Cm(2),
    Cm(2),
    Cm(18),
    Cm(2)
)

p = textbox.text_frame.paragraphs[0]
p.text = "Hello PowerPoint"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 70, 180)

# ===== 藍色框 =====
shape = slide.shapes.add_shape(
    MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
    Cm(2),
    Cm(5),
    Cm(18),
    Cm(2)
)

shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0, 70, 180)
shape.line.color.rgb = RGBColor(0, 70, 180)

tf = shape.text_frame
tf.text = "這是第一個 PowerPoint"

para = tf.paragraphs[0]
para.alignment = PP_ALIGN.CENTER
para.font.size = Pt(20)
para.font.bold = True
para.font.color.rgb = RGBColor(255, 255, 255)

# 存檔
prs.save("第一份PPT.pptx")

print("完成！")