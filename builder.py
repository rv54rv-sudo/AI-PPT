from pptx import Presentation
from pptx.util import Cm
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor

from components import 加入文字, 加入提醒框


def create_ppt():

    prs = Presentation()

    prs.slide_width = Cm(21)
    prs.slide_height = Cm(29.7)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 背景
    #slide.shapes.add_picture(
        #"圖片素材/poster.jpg",
        #0,
        #0,
        #width=prs.slide_width,
        #height=prs.slide_height
    #)



# 公司 Logo
slide.shapes.add_picture(
    "公司LOGO/logo.png",
    Cm(0.8),
    Cm(0.6),
    width=Cm(2.4),
    height=Cm(2.4)
)

    # 公司名稱
    加入文字(
        slide,
        "欣欣天然氣股份有限公司",
        3.6,
        0.9,
        10,
        0.8,
        字型大小=18,
        粗體=True,
        顏色=(0,80,170)
    )

    # 安全指引
    加入文字(
        slide,
        "安心用氣｜安全指引",
        15.5,
        1.1,
        4,
        0.6,
        字型大小=10,
        顏色=(80,80,80)
    )

    # 主標題
    加入文字(
        slide,
        "天然氣突然停止供氣？",
        1.2,
        3.6,
        16,
        1.5,
        字型大小=28,
        粗體=True,
        顏色=(0,70,180)
    )
加入提醒框(slide)
    prs.save("輸出檔案/Header_v1.pptx")

    print("Header 完成")