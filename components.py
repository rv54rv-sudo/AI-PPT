from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def 加入文字(
    slide,
    文字,
    x,
    y,
    寬,
    高,
    字型大小=18,
    粗體=False,
    顏色=(0,0,0),
    置中=False,
):

    textbox = slide.shapes.add_textbox(
        Cm(x),
        Cm(y),
        Cm(寬),
        Cm(高)
    )

    tf = textbox.text_frame
    p = tf.paragraphs[0]

    p.text = 文字
    p.font.size = Pt(字型大小)
    p.font.bold = 粗體
    p.font.color.rgb = RGBColor(*顏色)

    if 置中:
        p.alignment = PP_ALIGN.CENTER

    return textboxfrom pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

def 加入提醒框(slide):

    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Cm(0.9),
        Cm(5.0),
        Cm(7.4),
        Cm(1.8)
    )

    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(43,145,52)

    box.line.color.rgb = RGBColor(43,145,52)

    tf = box.text_frame

    p = tf.paragraphs[0]
    p.text = "請別擔心！"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255,255,255)

    p = tf.add_paragraph()
    p.text = "微電腦瓦斯表可能已啟動\n安全遮斷功能。"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255,255,255)